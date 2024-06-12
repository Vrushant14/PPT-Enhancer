import copy
from openai import OpenAI
from pptx import Presentation
from pptx.util import Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import os
from dotenv import load_dotenv
import pymongo
from myproject import settings
import openai
import requests
from pptx.util import Inches



load_dotenv()
api_key = os.getenv('api_key')

client = OpenAI(api_key=api_key)

client1 = pymongo.MongoClient(settings.MONGODB_URI)
db = client1[settings.MONGODB_NAME]
collection = db['enhanced_ppt']
collection1 = db['gen_images']

def extract_text_from_ppt(ppt_file):
    prs = Presentation(ppt_file)

    extracted_text_per_slide = []

    for slide in prs.slides:
        slide_text = ""

        for shape in slide.shapes:
            if hasattr(shape, "text"):
                slide_text += shape.text + "\n"

        extracted_text_per_slide.append(slide_text)

    return extracted_text_per_slide

def enhance_text_with_openai(text):

    engine = "gpt-3.5-turbo-1106"  # Replace with the engine you want to use

    response = client.chat.completions.create(model=engine,
                                              messages=[{"role": "user", "content": text}],
                                              max_tokens=150)

    enhanced_text = response.choices[0].message.content

    return enhanced_text


def SlideCopyFromPasteInto(copyFromPres, slideIndex,  pasteIntoPres):

    slide_to_copy = copyFromPres.slides[slideIndex]

    slide_layout = pasteIntoPres.slide_layouts.get_by_name("Blank") 
    
    new_slide = pasteIntoPres.slides.add_slide(slide_layout)

    imgDict = {}

   
    for shp in slide_to_copy.shapes:
        if 'Picture' in shp.name:
            with open(shp.name+'.jpg', 'wb') as f:
                f.write(shp.image.blob)

            imgDict[shp.name+'.jpg'] = [shp.left, shp.top, shp.width, shp.height]
        else:
            el = shp.element
            newel = copy.deepcopy(el)

            new_slide.shapes._spTree.insert_element_before(newel, 'p:extLst')
    
    for k, v in imgDict.items():
        new_slide.shapes.add_picture(k, v[0], v[1], v[2], v[3])
        os.remove(k)

    return new_slide 

def process_ppt(ppt_file_path, username):
  try:
    prs = Presentation(r"C:\Users\Vrushant\Desktop\pblfinal\PBL-master\accounts\template.pptx")
    extracted_text_per_slide = extract_text_from_ppt(ppt_file_path)
    enhanced_text_per_slide = [enhance_text_with_openai(text) for text in extracted_text_per_slide]
    for enhanced_text in enhanced_text_per_slide:
        new_slide = SlideCopyFromPasteInto(prs, 0, prs)  # Title and Content layout
        for shape in new_slide.shapes:
            if hasattr(shape, "text"):
                shape.text = enhanced_text
                for paragraph in shape.text_frame.paragraphs:
                            for run in paragraph.runs:
                                run.font.size = Pt(18)
        keywords = generate_keywords(enhanced_text)
        image_id = generate_image(keywords)

        if image_id:
            image_data_document = collection1.find_one({"_id": image_id})
            if image_data_document:
                image_data = image_data_document.get("data")
                if image_data:
                    with open("temp_image.jpg", "wb") as f:
                        f.write(image_data)
                    left = Inches(9.2)
                    top = Inches(1.5)
                    pic = new_slide.shapes.add_picture("temp_image.jpg", left, top, width=Inches(3.8), height=Inches(3.8))
                else:
                    print("Image data not found in the document.")
            else:
                print("Image document not found in the collection.")
        else:
            print("Image ID is None.")

    xml_slides = prs.slides._sldIdLst  
    slides = list(xml_slides)
    xml_slides.remove(slides[0]) 
    enhanced_ppt_path = "enhanced_presentation.pptx"
    prs.save(enhanced_ppt_path)

    
    return enhanced_ppt_path
  except Exception as e:
    print("An error occurred:", str(e))

def generate_keywords(text):
    engine = "gpt-3.5-turbo-1106" 
    response = openai.ChatCompletion.create(

        model=engine,

        messages=[

            {"role": "system", "content": "You are a keyword extractor."},

            {"role": "user", "content": f"Extract the most relevant keywords from the following text: '{text}'."},

        ],

    )

    keywords = response.choices[0].message.content.strip().split("\n")


    # Frame a sentence for image generative AI

    response = openai.ChatCompletion.create(

        model=engine,

        messages=[

            {"role": "system", "content": "You are a helpful assistant that frames sentences for image generative AI."},

            {"role": "user", "content": f"Frame a sentence for image generative AI using the following keywords: {', '.join(keywords)}. The sentence should be in the format: 'Generate an image of [KEYWORD]'. Make sure to include all keywords."},

        ],

    )

    sentence = response.choices[0].message.content.strip()
    
     # Replace with the engine you want to use
    # text = text + "give me the most important keyword, in such a sentence format, that I can use to generate image using Dall E"
    # response = client.chat.completions.create(model=engine,
    #                                           messages=[{"role": "user", "content": text}],
    #                                           max_tokens=150)

    # enhanced_text = response.choices[0].message.content
    
    return sentence


def generate_image(text):
    openai.api_key = api_key
    keywords = generate_keywords(text)
    print(keywords)
    
    response = openai.images.generate(
        prompt="generate a basic animated image for the following without any text in the picture: "+keywords,
        n=1,
        size="1024x1024"  
    )

    image_url = response.data[0].url
    
    image_response = requests.get(image_url)

    if image_response.status_code == 200:
        img_data = image_response.content
        if img_data:
            img = {'name': 'img.jpg', 'data': img_data}
            result = collection1.insert_one(img)
            img_id = result.inserted_id
            print("Image saved successfully to MongoDB.")
            return img_id
    else:
        print("Failed to download the image:", image_response.status_code)